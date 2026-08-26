#!/usr/bin/env python3
"""Build Track 2 figures, HTML deck, and PPTX.

Each figure uses a clear layout. Labels do not overlap.
All visible text uses Simplified Technical English.
"""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mp
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
DPI = 200
INK = "#1a1a1a"
MUTED = "#555555"
NAVY = "#1e3a5f"
BLUE = "#93c5fd"
YELLOW = "#fde68a"
PINK = "#fbcfe8"
GREEN = "#86efac"
RED = "#fca5a5"
CYAN = "#67e8f9"
CREAM = "#fef3c7"
LILAC = "#ddd6fe"


def _style():
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 11,
            "axes.titlesize": 14,
            "text.color": INK,
            "axes.labelcolor": INK,
            "axes.edgecolor": INK,
            "figure.facecolor": "white",
            "savefig.facecolor": "white",
            "savefig.bbox": "tight",
            "savefig.pad_inches": 0.18,
        }
    )


def fig_evidence(path: Path) -> None:
    """Two-allele evidence card."""
    rows = [
        ("Allele", "Allele 1", "Allele 2"),
        ("Locus", "chr15:40209701 T>G", "chr15:40220612 T>G"),
        ("HGVS", "c.2210T>G  p.Leu737Ter", "c.3006T>G  p.Asn1002Lys"),
        ("Class", "stop-gain (HIGH)", "missense, kinase domain"),
        ("ClinVar", "P/LP  VCV000533901", "not in ClinVar (new)"),
        ("gnomAD AF", "7.9e-05", "6.8e-07 (one allele)"),
        ("Prediction", "stop-gain. No missense score.", "SIFT 0.01. PP2 0.997. AM 0.9229"),
        ("Genotype", "het PASS  AD 21/25  GQ 99", "het PASS  AD 15/13  GQ 99"),
    ]
    fig, ax = plt.subplots(figsize=(13.3, 4.6))
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_title("Causal pair. Two BUB1B alleles.", pad=10, fontweight="bold")
    col_x = [0.02, 0.28, 0.64]
    col_w = [0.24, 0.34, 0.34]
    n = len(rows)
    row_h = 0.82 / n
    for i, (a, b, c) in enumerate(rows):
        y = 0.90 - (i + 1) * row_h
        bg = "#f3f4f6" if i % 2 else "white"
        if i == 0:
            bg = "#e8eef5"
        ax.add_patch(mp.Rectangle((0.01, y), 0.98, row_h, fc=bg, ec="#d1d5db", lw=0.6))
        weight = "bold" if i == 0 else "normal"
        ax.text(col_x[0] + 0.01, y + row_h / 2, a, va="center", ha="left", fontweight="bold", fontsize=11)
        ax.text(col_x[1] + 0.01, y + row_h / 2, b, va="center", ha="left", fontweight=weight, fontsize=11)
        ax.text(col_x[2] + 0.01, y + row_h / 2, c, va="center", ha="left", fontweight=weight, fontsize=11)
    fig.savefig(path, dpi=DPI)
    plt.close(fig)


def fig_mechanism(path: Path) -> None:
    """Genotype-to-drug flow. Three target groups."""
    fig, ax = plt.subplots(figsize=(13.3, 7.0))
    ax.set_xlim(0, 13.3)
    ax.set_ylim(0, 7.0)
    ax.axis("off")
    ax.set_title("From genotype to drug. Three target groups.", fontweight="bold", pad=8)

    def box(x, y, w, h, text, fc, fs=10.5):
        p = FancyBboxPatch(
            (x, y), w, h, boxstyle="round,pad=0.04,rounding_size=0.08",
            fc=fc, ec="#374151", lw=1.1,
        )
        ax.add_patch(p)
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=fs, wrap=True)
        return (x + w / 2, y + h / 2, x, y, w, h)

    def arrow(a, b):
        ax.annotate(
            "",
            xy=(b[0], b[1] + b[5] / 2 if b[1] < a[1] else b[1] - 0),
            xytext=(a[0], a[1] - a[5] / 2 if b[1] < a[1] else a[1] + a[5] / 2),
            arrowprops=dict(arrowstyle="-|>", color="#374151", lw=1.4),
        )

    def connect(src, dst, src_edge="bottom", dst_edge="top"):
        sx, sy, x, y, w, h = src
        dx, dy, x2, y2, w2, h2 = dst
        if src_edge == "bottom":
            p0 = (sx, y)
        elif src_edge == "right":
            p0 = (x + w, sy)
        else:
            p0 = (sx, y + h)
        if dst_edge == "top":
            p1 = (dx, y2 + h2)
        elif dst_edge == "left":
            p1 = (x2, dy)
        else:
            p1 = (dx, y2)
        ax.add_patch(
            FancyArrowPatch(
                p0, p1, arrowstyle="-|>", mutation_scale=12, lw=1.3, color="#374151",
            )
        )

    a = box(0.35, 5.55, 3.9, 1.05, "BUB1B compound heterozygote\np.Leu737Ter + p.Asn1002Lys", "#fecaca")
    b = box(4.70, 5.55, 3.9, 1.05, "BubR1 is below the SAC threshold.\nMCC fails.", "#fecaca")
    c = box(9.05, 5.55, 3.9, 1.05, "APC/C starts.\nCells enter anaphase too early.", "#fecaca")
    m = box(4.70, 3.75, 3.9, 1.05, "Mosaic variegated\naneuploidy (MVA)", "#c7d2fe")
    s1 = box(0.35, 1.95, 3.9, 1.15, "Embryonal rhabdomyosarcoma.\nWilms tumor risk.", "#bbf7d0")
    s2 = box(4.70, 1.95, 3.9, 1.15, "Stress in aneuploid cells.\nmTOR, cGAS-STING, SASP.", "#bbf7d0")
    s3 = box(9.05, 1.95, 3.9, 1.15, "Severe IUGR.\nFailure to thrive. Prematurity.", "#bbf7d0")
    t1 = box(0.35, 0.25, 3.9, 1.20, "Tier A\nadavosertib + irinotecan\ncitrate + thiazide", CREAM)
    t2 = box(4.70, 0.25, 3.9, 1.20, "Tier B\neverolimus\n(mTOR buffer)", CREAM)
    t3 = box(9.05, 0.25, 3.9, 1.20, "Tier C\nrestore BubR1\n(AAV / SIRT2 / peptide)", CREAM)

    connect(a, b, "right", "left")
    connect(b, c, "right", "left")
    connect(c, m, "bottom", "top")
    connect(m, s1, "bottom", "top")
    connect(m, s2, "bottom", "top")
    connect(m, s3, "bottom", "top")
    connect(s1, t1, "bottom", "top")
    connect(s2, t2, "bottom", "top")
    connect(s3, t3, "bottom", "top")
    fig.savefig(path, dpi=DPI)
    plt.close(fig)


def fig_domainmap(path: Path) -> None:
    """Two-track BubR1 domain map. Full length plus kinase zoom."""
    fig, (ax1, ax2) = plt.subplots(
        2, 1, figsize=(13.3, 6.2), gridspec_kw={"height_ratios": [1.0, 1.1], "hspace": 0.38}
    )
    length = 1050

    def hide(ax):
        ax.set_xticks([])
        ax.set_yticks([])
        for sp in ax.spines.values():
            sp.set_visible(False)

    def backbone(ax, x0, x1, y):
        ax.plot([x0, x1], [y, y], color="#9ca3af", lw=7, solid_capstyle="round", zorder=1)

    def domain(ax, x0, x1, y, h, fc, label, fs=9.5):
        ax.add_patch(mp.Rectangle((x0, y - h / 2), x1 - x0, h, fc=fc, ec="#374151", lw=0.8, zorder=2))
        if (x1 - x0) >= 80:
            ax.text((x0 + x1) / 2, y, label, ha="center", va="center", fontsize=fs, zorder=3)
        else:
            ax.text((x0 + x1) / 2, y + h / 2 + 0.07, label, ha="center", va="bottom", fontsize=8, zorder=3)

    ax1.set_xlim(-20, 1100)
    ax1.set_ylim(-0.12, 1.28)
    hide(ax1)
    ax1.set_title("BubR1 protein (O60566). Length: 1050 amino acids.", fontweight="bold", loc="left")
    y = 0.52
    backbone(ax1, 1, length, y)
    domain(ax1, 1, 226, y, 0.26, BLUE, "TPR  1–226")
    domain(ax1, 152, 185, y, 0.14, YELLOW, "")
    ax1.text(168, y - 0.26, "KNL1 152–185", ha="center", va="top", fontsize=8)
    domain(ax1, 224, 232, y, 0.26, PINK, "")
    ax1.text(260, y + 0.20, "D-box", ha="center", va="bottom", fontsize=8)
    domain(ax1, 668, 675, y, 0.26, GREEN, "")
    ax1.text(672, y + 0.20, "B56", ha="center", va="bottom", fontsize=8)
    domain(ax1, 766, 1050, y, 0.26, RED, "Kinase domain  766–1050")
    for x, lab in ((20, "KEN"), (470, "ABBA")):
        ax1.plot([x, x], [y - 0.18, y + 0.18], color="#4b5563", lw=1.5, zorder=4)
        ax1.text(x, y - 0.22, lab, ha="center", va="top", fontsize=8, color="#374151")
    ax1.annotate(
        "Allele 1 stops here\np.Leu737Ter",
        xy=(737, y + 0.14), xytext=(560, 1.12),
        ha="center", fontsize=9, color="#991b1b",
        arrowprops=dict(arrowstyle="-|>", color="#991b1b", lw=1.4),
    )
    ax1.annotate(
        "Allele 2 changes this residue\np.Asn1002Lys",
        xy=(1002, y + 0.14), xytext=(920, 1.12),
        ha="center", fontsize=9, color="#9a3412",
        arrowprops=dict(arrowstyle="-|>", color="#9a3412", lw=1.4),
    )
    for tick in (1, 200, 400, 600, 800, 1000, 1050):
        ax1.plot([tick, tick], [0.22, 0.26], color="#9ca3af", lw=0.8)
        ax1.text(tick, 0.18, str(tick), ha="center", va="top", fontsize=7, color=MUTED)
    ax1.text(525, 0.02, "Residue number", ha="center", fontsize=8, color=MUTED)

    ax2.set_xlim(620, 1080)
    ax2.set_ylim(-0.28, 1.32)
    hide(ax2)
    ax2.set_title("Zoom. Residues 640–1050. Scaffold stays. Kinase is lost or weak.", fontweight="bold", loc="left")
    y2 = 0.70
    backbone(ax2, 640, 1050, y2)
    domain(ax2, 668, 675, y2, 0.24, GREEN, "")
    ax2.text(671.5, y2 + 0.20, "B56 motif 668–675", ha="center", va="bottom", fontsize=8)
    domain(ax2, 766, 1050, y2, 0.28, RED, "Kinase domain")
    ax2.plot([737, 737], [y2 - 0.20, y2 + 0.22], color="#991b1b", ls="--", lw=1.2)
    ax2.text(737, 1.12, "stop at 737", ha="center", fontsize=8, color="#991b1b")
    ax2.plot(1002, y2 + 0.16, "v", color="#9a3412", ms=8)
    ax2.text(1002, 1.12, "N1002K", ha="center", fontsize=8, color="#9a3412")
    for x, lab in ((772, "K772\nVAIK"), (795, "F795\nATP"), (882, "D882\ncatalytic")):
        ax2.plot([x, x], [y2 - 0.20, y2 - 0.15], color="#7f1d1d", lw=1.3)
        ax2.text(x, y2 - 0.23, lab, ha="center", va="top", fontsize=7.5, color="#7f1d1d")
    ax2.add_patch(mp.Rectangle((668, 0.18), 8, 0.10, fc=GREEN, ec="#374151", lw=0.6))
    ax2.text(680, 0.23, "PDB 5JJA. PP2A-B56 contact. Allele 1 keeps this motif.", va="center", fontsize=8)
    ax2.add_patch(mp.Rectangle((640, -0.04), 36, 0.10, fc=CYAN, ec="#374151", lw=0.6))
    ax2.text(680, 0.01, "PDB 6TLJ / 5KHU. APC/C-MCC scaffold (residues 19–499) is intact on allele 1.", va="center", fontsize=8)
    ax2.text(
        850, -0.20,
        "AlphaFold v6. Kinase mean pLDDT 82. N1002 pLDDT 91. Distance N1002 to D882 = 19.8 Å.",
        ha="center", fontsize=8, color=MUTED, style="italic",
    )
    fig.savefig(path, dpi=DPI)
    plt.close(fig)


def fig_exomiser(path: Path) -> None:
    """Exomiser top-gene bar chart from the genome-wide TSV."""
    tsv = HERE.parents[1] / "track1" / "analysis" / "exomiser" / "PROBAND01_genomewide.genes.tsv"
    rows: list[tuple[str, float]] = []
    seen: set[str] = set()
    for line in tsv.read_text().splitlines():
        if not line or line.startswith("#"):
            continue
        parts = line.split("\t")
        gene, moi, score = parts[2], parts[4], float(parts[6])
        key = f"{gene} {moi}"
        if key in seen:
            continue
        seen.add(key)
        rows.append((f"{gene} ({moi})", score))
        if len(rows) >= 8:
            break
    fig, ax = plt.subplots(figsize=(13.3, 6.4))
    labels = [g for g, _ in rows][::-1]
    vals = [v for _, v in rows][::-1]
    colors = ["#b91c1c" if lab.startswith("BUB1B") else "#64748b" for lab in labels]
    bars = ax.barh(labels, vals, color=colors, edgecolor="#111827", height=0.62)
    ax.set_xlabel("Exomiser combined score. No gene panel. Runtime 94 s.")
    ax.set_title("Exomiser ranks genes from the full genome.", fontweight="bold")
    ax.set_xlim(0, max(vals) * 1.18)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for bar, val in zip(bars, vals):
        ax.text(val + 0.008, bar.get_y() + bar.get_height() / 2, f"{val:.3f}", va="center", fontsize=10)
    fig.savefig(path, dpi=DPI)
    plt.close(fig)


NAVY_RGB = RGBColor(0x1E, 0x3A, 0x5F)
INK_RGB = RGBColor(0x1A, 0x1A, 0x1A)
MUTED_RGB = RGBColor(0x4B, 0x55, 0x63)
ACCENT = RGBColor(0x9A, 0x34, 0x12)


def _set_run(run, text, size=20, bold=False, color=INK_RGB):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_title(slide, text, top=0.28):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_run(p.add_run(), text, size=28, bold=True, color=NAVY_RGB)
    return box


def _add_bullets(slide, items, top=1.15, size=20):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(12)
        _set_run(p.add_run(), item, size=size)
    return box


def _add_caption(slide, text, top=6.85):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    _set_run(p.add_run(), text, size=13, color=MUTED_RGB)
    return box


def _add_refs(slide, items, top=1.05, size=12):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(6.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        _set_run(p.add_run(), item, size=size, color=INK_RGB)
    return box


def _tier_block(slide, y, title, body, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(1.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = color
    shape.line.width = Emu(19050)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_run(p.add_run(), title, size=18, bold=True, color=color)
    p2 = tf.add_paragraph()
    _set_run(p2.add_run(), body, size=16)


def build_pptx(out: Path, figs: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    _add_title(s, "A real child. A checkpoint that does not work. A three-tier plan.")
    _add_bullets(
        s,
        [
            "Mosaic Variegated Aneuploidy (MVA). Cells lose and gain chromosomes.",
            "Embryonal rhabdomyosarcoma. Severe growth failure. Nephrocalcinosis.",
            "Team bigbag. Rare Disease, Real Kid: MVA Hackathon 2026.",
        ],
        top=1.6,
        size=22,
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "The cause. Confirmed.")
    s.shapes.add_picture(str(figs["evidence"]), Inches(0.7), Inches(1.15), width=Inches(11.9))
    _add_caption(s, "Track 1 score: 100.0 rank points. F-max 1.000. Full match.")

    s = prs.slides.add_slide(blank)
    _add_title(s, "The constraint")
    _add_bullets(
        s,
        [
            "No approved drug increases BubR1.",
            "We must not decrease checkpoint function more. The checkpoint already does not work.",
            "We map three target groups: the tumor, the stress, and the protein.",
        ],
        top=1.8,
        size=24,
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "Mechanism. From genotype to drug.")
    s.shapes.add_picture(str(figs["mechanism"]), Inches(1.15), Inches(1.08), width=Inches(11.0))

    s = prs.slides.add_slide(blank)
    _add_title(s, "Tier A. Treat the signs that we see today.")
    _tier_block(
        s, 1.25, "Adavosertib plus irinotecan",
        "A pediatric trial sets the dose. NCT02095132. https://clinicaltrials.gov/study/NCT02095132",
        RGBColor(0x16, 0x65, 0x34),
    )
    _tier_block(
        s, 2.85, "Potassium citrate plus a thiazide",
        "Standard care for nephrocalcinosis. Weigert and Hoppe. https://pubmed.ncbi.nlm.nih.gov/29707529/",
        RGBColor(0x16, 0x65, 0x34),
    )
    _tier_block(
        s, 4.45, "Alisertib",
        "Alternate mitotic drug. COG ADVL0921. https://pubmed.ncbi.nlm.nih.gov/30777875/",
        RGBColor(0x16, 0x65, 0x34),
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "Tier B. Decrease the stress.")
    _tier_block(
        s, 1.25, "Aneuploid cells increase mTOR.",
        "These cells release inflammatory signals.",
        RGBColor(0x1D, 0x4E, 0x89),
    )
    _tier_block(
        s, 2.85, "Everolimus decreases that signal.",
        "Pediatric safety data are large. EXIST-1 long-term. https://pubmed.ncbi.nlm.nih.gov/27351628/",
        RGBColor(0x1D, 0x4E, 0x89),
    )
    _tier_block(
        s, 4.45, "Honest claim",
        "This drug decreases damage. This drug does not repair the checkpoint.",
        RGBColor(0x1D, 0x4E, 0x89),
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "Tier C. Restore the protein.")
    _tier_block(
        s, 1.25, "Mouse data (Baker 2013)",
        "More BubR1 repairs the checkpoint. https://pubmed.ncbi.nlm.nih.gov/23242215/",
        RGBColor(0x9A, 0x34, 0x12),
    )
    _tier_block(
        s, 2.85, "This child still makes weakened BubR1 from the missense allele.",
        "A small increase in BubR1 level can restore checkpoint function.",
        RGBColor(0x9A, 0x34, 0x12),
    )
    _tier_block(
        s, 4.45, "Three routes. One goal.",
        "AAV gene augmentation. SIRT2-axis stabilization. Peptide rescue.",
        RGBColor(0x9A, 0x34, 0x12),
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "Both alleles change the kinase domain. The scaffold stays.")
    s.shapes.add_picture(str(figs["domain"]), Inches(0.45), Inches(1.00), width=Inches(12.4))

    s = prs.slides.add_slide(blank)
    _add_title(s, "Rigor. Three independent methods agree.")
    s.shapes.add_picture(str(figs["exomiser"]), Inches(0.7), Inches(1.05), width=Inches(11.9))
    _add_caption(s, "Panel scan. Exomiser on the full genome (94 s, no panel). Clinical databases. All point to BUB1B.")

    s = prs.slides.add_slide(blank)
    _add_title(s, "Scale. A laptop runs the pipeline.")
    _add_bullets(
        s,
        [
            "Open Targets. DGIdb. ChEMBL. everycure/matrix-scores (39.5 million pairs). All public APIs.",
            "Full pipeline: about 6 hours and 1 USD of electricity. Genome-wide check: 94 seconds.",
            "The method applies to other recessive chromosomal-instability disorders.",
            "Code and evidence: github.com/bigbag/mva-hackathon-2026",
        ],
        top=1.7,
        size=22,
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "The request")
    _add_bullets(
        s,
        [
            "Give standard symptom care now.",
            "Use trial-controlled cancer options when needed.",
            "Start a research program on BubR1 restoration.",
            "This family shared the genome of their child. We give them a clear plan.",
        ],
        top=1.7,
        size=24,
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "Works Cited (MLA)")
    _add_refs(
        s,
        [
            "Baker et al. Nat Cell Biol 2013. https://doi.org/10.1038/ncb2643  https://pubmed.ncbi.nlm.nih.gov/23242215/",
            "Franz et al. Lancet 2013 EXIST-1. https://doi.org/10.1016/S0140-6736(12)61134-9  https://pubmed.ncbi.nlm.nih.gov/23158522/",
            "Franz et al. PLOS ONE 2016 EXIST-1 long-term. https://doi.org/10.1371/journal.pone.0158476  https://pubmed.ncbi.nlm.nih.gov/27351628/",
            "Hanks et al. Nat Genet 2004. https://doi.org/10.1038/ng1449  https://pubmed.ncbi.nlm.nih.gov/15475955/",
            "Malumbres and Villarroya-Beltri. Nat Rev Genet 2024. https://doi.org/10.1038/s41576-024-00762-6  https://pubmed.ncbi.nlm.nih.gov/39169218/",
            "Mossé et al. Clin Cancer Res 2019 ADVL0921. https://doi.org/10.1158/1078-0432.CCR-18-2675  https://pubmed.ncbi.nlm.nih.gov/30777875/",
            "National Cancer Institute. NCT02095132. https://clinicaltrials.gov/study/NCT02095132",
            "Sage Bionetworks. Hackathon. https://sagebio-rare-disease-real-kid-mva-hackathon-2026.hf.space/",
            "Sage Bionetworks. Data. https://huggingface.co/datasets/SageBio/mva-hackathon-2026-data",
        ],
        size=13,
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "Works Cited (MLA), continued")
    _add_refs(
        s,
        [
            "Weigert and Hoppe. Front Pediatr 2018. https://doi.org/10.3389/fped.2018.00098  https://pubmed.ncbi.nlm.nih.gov/29707529/",
            "Yost et al. Nat Genet 2017. https://doi.org/10.1038/ng.3883  https://pubmed.ncbi.nlm.nih.gov/28553959/",
            "Zhang et al. Nat Aging 2023. https://doi.org/10.1038/s43587-023-00361-w  https://pubmed.ncbi.nlm.nih.gov/37118121/",
            "Villarroya-Beltri et al. Sci Adv 2022. https://doi.org/10.1126/sciadv.abq5914  https://pubmed.ncbi.nlm.nih.gov/36322655/",
            "Cole et al. Cancer 2023. https://doi.org/10.1002/cncr.34786  https://pubmed.ncbi.nlm.nih.gov/37081608/",
            "Open Targets. BUB1B. https://platform.opentargets.org/target/ENSG00000156970",
            "ClinVar. VCV000533901. https://www.ncbi.nlm.nih.gov/clinvar/variation/533901/",
            "bigbag. Code. https://github.com/bigbag/mva-hackathon-2026",
        ],
        size=13,
    )

    s = prs.slides.add_slide(blank)
    _add_title(s, "Thank you")
    _add_bullets(
        s,
        [
            "Team bigbag. https://github.com/bigbag/mva-hackathon-2026",
            "Hackathon: https://sagebio-rare-disease-real-kid-mva-hackathon-2026.hf.space/",
            "Data: https://huggingface.co/datasets/SageBio/mva-hackathon-2026-data",
            "Organizers: Sage Bionetworks, MVA Society, Hugging Face, BEACON. Sponsors: AWS and Anthropic.",
            "Above all, the family.",
        ],
        top=1.5,
        size=18,
    )

    prs.save(out)


def build_html(out: Path) -> None:
    html = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>MVA Hackathon 2026. bigbag pitch deck</title>
<style>
  body { margin:0; font-family: Calibri, 'Segoe UI', system-ui, sans-serif; background:#111; color:#eee; }
  .slide { width:1280px; height:720px; margin:20px auto; background:#f8fafc; color:#1a1a1a;
           border-radius:8px; padding:48px 64px; box-sizing:border-box; display:flex; flex-direction:column; }
  h1 { font-size:34px; margin:0 0 22px; color:#1e3a5f; }
  li { font-size:24px; margin:12px 0; line-height:1.35; }
  img { max-width:100%; max-height:540px; margin:auto; }
  .cap { color:#4b5563; font-size:16px; text-align:center; margin-top:10px; }
  .tier { border-left:6px solid; padding:14px 18px; margin:14px 0; background:#fff; border-radius:6px; }
  .a { border-color:#166534; } .b { border-color:#1d4e89; } .c { border-color:#9a3412; }
  em { color:#9a3412; font-style:normal; font-weight:600; }
</style>
</head>
<body>

<section class="slide"><h1>A real child. A checkpoint that does not work. A three-tier plan.</h1>
<ul>
<li>Mosaic Variegated Aneuploidy (MVA). Cells lose and gain chromosomes.</li>
<li>Embryonal rhabdomyosarcoma. Severe growth failure. Nephrocalcinosis.</li>
<li>Team bigbag. Rare Disease, Real Kid: MVA Hackathon 2026.</li>
</ul></section>

<section class="slide"><h1>The cause. Confirmed.</h1>
<img src="fig_evidence.png" alt="Evidence card for the two BUB1B alleles">
<div class="cap">Track 1 score: 100.0 rank points. F-max 1.000. Full match.</div></section>

<section class="slide"><h1>The constraint</h1>
<ul>
<li>No approved drug increases BubR1.</li>
<li>We must not decrease checkpoint function more. The checkpoint already does not work.</li>
<li>We map three target groups: <em>the tumor, the stress, and the protein.</em></li>
</ul></section>

<section class="slide"><h1>Mechanism. From genotype to drug.</h1>
<img src="fig_mechanism.png" alt="Flow from genotype to three drug tiers"></section>

<section class="slide"><h1>Tier A. Treat the signs that we see today.</h1>
<div class="tier a"><b>Adavosertib plus irinotecan.</b> A pediatric trial sets the dose (<a href="https://clinicaltrials.gov/study/NCT02095132">NCT02095132</a>). The trial has a rhabdomyosarcoma arm.</div>
<div class="tier a"><b>Potassium citrate plus a thiazide.</b> Standard care for nephrocalcinosis (<a href="https://pubmed.ncbi.nlm.nih.gov/29707529/">Weigert and Hoppe</a>). Safe in children. Protects bone.</div>
<div class="tier a"><b>Alisertib.</b> Alternate mitotic drug. COG phase 2 trial <a href="https://pubmed.ncbi.nlm.nih.gov/30777875/">ADVL0921</a>.</div></section>

<section class="slide"><h1>Tier B. Decrease the stress.</h1>
<div class="tier b">Aneuploid cells increase mTOR. These cells release inflammatory signals.</div>
<div class="tier b"><b>Everolimus</b> decreases that signal. Pediatric safety data are large (<a href="https://pubmed.ncbi.nlm.nih.gov/27351628/">EXIST-1</a>, transplant).</div>
<div class="tier b">Honest claim: <em>this drug decreases damage. This drug does not repair the checkpoint.</em></div></section>

<section class="slide"><h1>Tier C. Restore the protein.</h1>
<div class="tier c">In mice, more BubR1 repairs the checkpoint and extends lifespan (<a href="https://pubmed.ncbi.nlm.nih.gov/23242215/">Baker 2013</a>).</div>
<div class="tier c">This child still makes weakened BubR1 from the missense allele.</div>
<div class="tier c">Three routes: <b>AAV augmentation. SIRT2-axis stabilization. Peptide rescue.</b></div></section>

<section class="slide"><h1>Both alleles change the kinase domain. The scaffold stays.</h1>
<img src="fig_domainmap.png" alt="BubR1 domain map with both mutations"></section>

<section class="slide"><h1>Rigor. Three independent methods agree.</h1>
<img src="fig_exomiser.png" alt="Exomiser ranks BUB1B first">
<div class="cap">Panel scan. Exomiser on the full genome (94 s, no panel). Clinical databases. All point to BUB1B.</div></section>

<section class="slide"><h1>Scale. A laptop runs the pipeline.</h1>
<ul>
<li>Open Targets. DGIdb. ChEMBL. everycure/matrix-scores (39.5 million pairs). All public APIs.</li>
<li>Full pipeline: about 6 hours and 1 USD of electricity. Genome-wide check: 94 seconds.</li>
<li>The method applies to other recessive chromosomal-instability disorders.</li>
<li>Code and evidence: github.com/bigbag/mva-hackathon-2026</li>
</ul></section>

<section class="slide"><h1>The request</h1>
<ul>
<li>Give standard symptom care now.</li>
<li>Use trial-controlled cancer options when needed.</li>
<li>Start a research program on BubR1 restoration.</li>
<li>This family shared the genome of their child. We give them a clear plan.</li>
</ul></section>

<section class="slide"><h1>Works Cited (MLA)</h1>
<ul style="font-size:16px;line-height:1.35">
<li>Baker et al. <em>Nat Cell Biol</em> 2013. <a href="https://pubmed.ncbi.nlm.nih.gov/23242215/">PubMed 23242215</a> · <a href="https://doi.org/10.1038/ncb2643">doi:10.1038/ncb2643</a></li>
<li>Franz et al. EXIST-1. <em>Lancet</em> 2013. <a href="https://pubmed.ncbi.nlm.nih.gov/23158522/">PubMed 23158522</a></li>
<li>Franz et al. EXIST-1 long-term. <em>PLOS ONE</em> 2016. <a href="https://pubmed.ncbi.nlm.nih.gov/27351628/">PubMed 27351628</a></li>
<li>Hanks et al. <em>Nat Genet</em> 2004. <a href="https://pubmed.ncbi.nlm.nih.gov/15475955/">PubMed 15475955</a></li>
<li>Malumbres and Villarroya-Beltri. <em>Nat Rev Genet</em> 2024. <a href="https://pubmed.ncbi.nlm.nih.gov/39169218/">PubMed 39169218</a></li>
<li>Mossé et al. ADVL0921. <em>Clin Cancer Res</em> 2019. <a href="https://pubmed.ncbi.nlm.nih.gov/30777875/">PubMed 30777875</a></li>
<li>NCT02095132. <a href="https://clinicaltrials.gov/study/NCT02095132">clinicaltrials.gov/study/NCT02095132</a></li>
<li>Hackathon. <a href="https://sagebio-rare-disease-real-kid-mva-hackathon-2026.hf.space/">HF Space</a></li>
<li>Data. <a href="https://huggingface.co/datasets/SageBio/mva-hackathon-2026-data">HF Dataset</a></li>
</ul></section>

<section class="slide"><h1>Works Cited (MLA), continued</h1>
<ul style="font-size:16px;line-height:1.35">
<li>Weigert and Hoppe. <em>Front Pediatr</em> 2018. <a href="https://pubmed.ncbi.nlm.nih.gov/29707529/">PubMed 29707529</a></li>
<li>Yost et al. <em>Nat Genet</em> 2017. <a href="https://pubmed.ncbi.nlm.nih.gov/28553959/">PubMed 28553959</a></li>
<li>Zhang et al. <em>Nat Aging</em> 2023. <a href="https://pubmed.ncbi.nlm.nih.gov/37118121/">PubMed 37118121</a></li>
<li>Villarroya-Beltri et al. <em>Sci Adv</em> 2022. <a href="https://pubmed.ncbi.nlm.nih.gov/36322655/">PubMed 36322655</a></li>
<li>Cole et al. <em>Cancer</em> 2023. <a href="https://pubmed.ncbi.nlm.nih.gov/37081608/">PubMed 37081608</a></li>
<li>Open Targets BUB1B. <a href="https://platform.opentargets.org/target/ENSG00000156970">ENSG00000156970</a></li>
<li>ClinVar. <a href="https://www.ncbi.nlm.nih.gov/clinvar/variation/533901/">VCV000533901</a></li>
<li>Code. <a href="https://github.com/bigbag/mva-hackathon-2026">github.com/bigbag/mva-hackathon-2026</a></li>
</ul></section>

<section class="slide"><h1>Thank you</h1>
<ul>
<li>Team bigbag. <a href="https://github.com/bigbag/mva-hackathon-2026">github.com/bigbag/mva-hackathon-2026</a></li>
<li>Hackathon: <a href="https://sagebio-rare-disease-real-kid-mva-hackathon-2026.hf.space/">sagebio-rare-disease-real-kid-mva-hackathon-2026.hf.space</a></li>
<li>Data: <a href="https://huggingface.co/datasets/SageBio/mva-hackathon-2026-data">huggingface.co/datasets/SageBio/mva-hackathon-2026-data</a></li>
<li>Organizers: Sage Bionetworks, MVA Society, Hugging Face, BEACON. Sponsors: AWS and Anthropic.</li>
<li>Above all, <em>the family</em>.</li>
</ul></section>
</body>
</html>
"""
    out.write_text(html)


def main() -> None:
    _style()
    figs = {
        "evidence": HERE / "fig_evidence.png",
        "mechanism": HERE / "fig_mechanism.png",
        "domain": HERE / "fig_domainmap.png",
        "exomiser": HERE / "fig_exomiser.png",
    }
    fig_evidence(figs["evidence"])
    fig_mechanism(figs["mechanism"])
    fig_domainmap(figs["domain"])
    fig_exomiser(figs["exomiser"])
    build_pptx(HERE / "pitch_deck.pptx", figs)
    build_html(HERE / "deck.html")
    print("wrote figures, pitch_deck.pptx, deck.html")


if __name__ == "__main__":
    main()
